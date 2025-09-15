import io
import time
import traceback
import mammoth
from PIL import Image
from docx import Document
import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from .html_converter import htmlConverter
from .utils.docx.pre_process import pre_process_docx
from ..utils.utils import extract_markdown_images,base64_to_pillow


class docxConverter(object):
    def __init__(self):
        self.html_converter = htmlConverter()
    
    def invoke_docx(self,file,ocrmodel, ocr_mode = "base", **kwargs):
        try:
            tt = time.time()
            content = []
            doc = Document(file)
            for eid,element in enumerate(doc.element.body):
                if element.tag.endswith('p'):
                    text = element.text
                    pics = element.xpath('.//pic:pic')
                    if text:
                        content.append(('text', text, eid))
                    elif pics:
                        for pid,pic in enumerate(pics):
                            blip = pic.find('.//a:blip', namespaces=pic.nsmap)
                            if blip is not None:
                                rId = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                                image_part = doc.part.related_parts[rId]
                                image_data = image_part.blob
                                image = Image.open(io.BytesIO(image_data))
                                content.append(('image', image, f"Img_{eid}_{pid}.png"))                    

                elif element.tag.endswith('tbl'):
                    # 处理表格
                    table = element
                    table_data = ""
                    for idx,row in enumerate(table.xpath('.//w:tr')):
                        row_data = []
                        for cell in row.xpath('.//w:tc'):
                            cell_text = ''.join(cell.itertext()).strip()
                            row_data.append(cell_text)
                        row_text = "|" + "|".join(row_data) + "|" + "\n"
                        if idx == 0:
                            row_text += "|" + "|".join(["-------" for x in row_data]) + "|" + "\n"
                        table_data += row_text
                    content.append(('table', table_data, eid))
            res = []
            full_text = ""
            images = {}
            for (content_type,content_text,content_name) in content:
                if content_type == 'text':
                    if full_text.endswith('\n'):
                        full_text += f'{content_text}\n'
                    else:
                        full_text += f'\n{content_text}\n'
                elif content_type == 'table':
                    if full_text.endswith('\n'):
                        full_text += f'{content_text}\n'
                    else:
                        full_text += f'\n{content_text}\n'
                elif content_type == 'image':
                    pil_image = content_text
                    # pil_image = resize_image_if_need(pil_image)
                    image_content = ocrmodel.invoke(pil_image, mode = ocr_mode, **kwargs)
                    image_content = image_content.to_markdown()
                    images[content_name] = {
                        "image": content,
                        "content": image_content
                    }
                    if image_content:
                        if full_text.endswith('\n'):
                            full_text += f"{image_content}\n"
                        else:
                            full_text += f"\n{image_content}\n"                                  

            res.append({
                "type": "docx",
                "content": full_text,
                "images": images,
                "time_elapse": time.time() - tt
            })
        except:
            traceback.print_exc()
            res = []
        finally:
            return res
    
    def invoke_docx_v2(self,file,ocrmodel, ocr_mode = "base", **kwargs):
        try:
            tt = time.time()
            res = []
            full_text = ""
            images = {}                                
            with open(file, "rb") as file_stream:
                pre_process_stream = pre_process_docx(file_stream)
                result = mammoth.convert_to_html(pre_process_stream)
            full_text = self.html_converter.invoke_string(result.value, **kwargs)
            full_text = full_text[0]["content"]
            image_matches = extract_markdown_images(full_text)
            replace_image_list = []
            for match in image_matches:
                try:
                    start_idx,end_idx = match.span()
                    full_match_text = match.group(0)
                    alt_text = match.group(1)
                    image_data = match.group(2)
                    image_data = image_data.split("base64,")[-1]
                    image_data = base64_to_pillow(image_data)
                    # image_data = resize_image_if_need(image_data)
                    image_content = ocrmodel.invoke(image_data, mode = ocr_mode, **kwargs)
                    image_content = image_content.to_markdown()   
                    image_content = f"\n{image_content}\n"
                    replace_image_list.append([
                        [start_idx,end_idx],
                        full_match_text,
                        image_content
                    ])            
                except:
                    traceback.print_exc()
                    continue
            replace_full_text = full_text
            # print(replace_image_list)
            try:
                # 先按 start 排序，并反向遍历（从后往前替换）
                replace_image_list = sorted(replace_image_list, key=lambda x: x[0][0], reverse=True)
                for box,ori_text,src_text in replace_image_list:
                    start_idx,end_idx = box
                    if start_idx > end_idx:
                        continue
                    # 替换 [start, end+1)，因为 Python 切片是前闭后开
                    replace_full_text = replace_full_text[:start_idx] + src_text + replace_full_text[end_idx + 1:]                    
            except:
                traceback.print_exc()
                replace_full_text = full_text           
            
            res.append({
                "type": "docx",
                "content": replace_full_text,
                "images": images,
                "time_elapse": time.time() - tt
            })
        except:
            traceback.print_exc()
            res = []
        finally:
            return res
                

class pptxConverter(object):
    def __init__(self):
        pass
    
    def invoke_pptx(self,file,ocrmodel, ocr_mode = "base", **kwargs):
        try:
            tt = time.time()
            content = []
            presentation = Presentation(file)
            for idx,slide in enumerate(presentation.slides):
                for sid,shape in enumerate(slide.shapes):
                    if shape.has_text_frame:
                        # 处理文本框
                        text = shape.text.strip()
                        if text:
                            content.append(('text', text, f"{idx}-{sid}"))
                    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        # 处理表格
                        table_data = ""
                        table = shape.table
                        for tid,row in enumerate(table.rows):
                            row_data = [cell.text.strip() for cell in row.cells]
                            row_text = "|" + "|".join(row_data) + "|" + "\n"
                            if tid == 0:
                                row_text += "|" + "|".join(["-------" for x in row_data]) + "|" + "\n"
                            table_data += row_text
                        content.append(('table', table_data, f"{idx}-{sid}"))
                    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        # 处理图片
                        image_stream = io.BytesIO(shape.image.blob)
                        image = Image.open(image_stream)
                        content.append(('image', image, f"Img_{idx}_{sid}.png"))
            res = []
            full_text = ""
            images = {}
            for (content_type,content_text,content_name) in content:
                if content_type == 'text':
                    if full_text.endswith('\n'):
                        full_text += f'{content_text}\n'
                    else:
                        full_text += f'\n{content_text}\n'
                elif content_type == 'table':
                    if full_text.endswith('\n'):
                        full_text += f'{content_text}\n'
                    else:
                        full_text += f'\n{content_text}\n'
                elif content_type == 'image':
                    pil_image = content_text
                    image_content = ocrmodel.invoke(pil_image, mode = ocr_mode, **kwargs)
                    image_content = image_content.to_markdown()
                    images[content_name] = {
                        "image": content,
                        "content": image_content
                    }
                    if image_content:
                        if full_text.endswith('\n'):
                            full_text += f"{image_content}\n"
                        else:
                            full_text += f"\n{image_content}\n"                                       

            res.append({
                "type": "pptx",
                "content": full_text,
                "images": images,
                "time_elapse": time.time() - tt
            })
        except:
            traceback.print_exc()
            res = []
        finally:
            return res
        

class tableConverter(object):
    def __init__(self):
        pass
    
    def invoke_csv(self,file,chunk_size = None,encoding = 'utf-8'):
        try:
            res = []
            if chunk_size:
                chunk_size = int(chunk_size)
                dfset = pd.read_csv(file,chunksize = chunk_size,encoding=encoding)
                for idx,df in enumerate(dfset):
                    df = df.dropna(how='all')
                    df = df.fillna("")
                    idx = idx + 1
                    res.append({
                        "sheet_name": "",
                        "chunk_id": f"{idx}",
                        "chunk_size": f"{chunk_size}",
                        "type": "csv",
                        "content": df.to_markdown(index = False)
                    })
            else:
                df = pd.read_csv(file,encoding=encoding)
                df = df.dropna(how='all')
                df = df.fillna("")
                res.append({
                    "sheet_name": "",
                    "chunk_id": "1",
                    "chunk_size": None,
                    "type": "csv",
                    "content": df.to_markdown(index = False)                
                })
        except:
            traceback.print_exc()
            res = []
        finally:
            return res
    
    def invoke_excel(self,file,chunk_size = None):
        try:
            res = []
            if chunk_size:
                chunk_size = int(chunk_size)
                dfset = pd.read_excel(file,sheet_name=None)
                for k in dfset.keys():
                    df = dfset[k]
                    df = df.dropna(how='all')
                    df = df.fillna("")
                    for idx,i in enumerate(range(0, len(df), chunk_size)):
                        idx = idx + 1
                        chunk = df.iloc[i:i + chunk_size]
                        res.append({
                            "sheet_name": f"{k}",
                            "chunk_id": f"{idx}",
                            "chunk_size": chunk_size,
                            "type": "excel",
                            "content": chunk.to_markdown(index = False)                     
                        })                                                
            else:
                dfset = pd.read_excel(file,sheet_name=None)
                for k in dfset.keys():
                    df = dfset[k]
                    df = df.dropna(how='all')
                    df = df.fillna("")
                    res.append({
                        "sheet_name": f"{k}",
                        "chunk_id": "1",
                        "chunk_size": None,
                        "type": "excel",
                        "content": df.to_markdown(index = False)                     
                    })
        except:
            traceback.print_exc()
            res = []
        finally:
            return res